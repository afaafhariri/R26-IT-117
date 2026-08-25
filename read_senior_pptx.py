from pptx import Presentation
from pptx.util import Inches, Pt

PPTX_PATH = r"C:\Users\Asus\Desktop\Research\PP1\25-26J-336_Empowering Students with Predictive Intelligence  A Career Guidance System for Sri Lanka.pptx"

prs = Presentation(PPTX_PATH)
print(f"Total slides: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width.inches:.1f} x {prs.slide_height.inches:.1f} inches")

for i, slide in enumerate(prs.slides):
    print(f"\n=== SLIDE {i+1} ===")
    print(f"Layout: {slide.slide_layout.name}")
    bg = slide.background.fill
    try:
        print(f"BG fill type: {bg.type}")
        if bg.fore_color:
            print(f"BG color: {bg.fore_color.rgb}")
    except:
        pass
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            tf = shape.text_frame
            for para in tf.paragraphs:
                line = para.text.strip()
                if not line:
                    continue
                sz = None
                bold = False
                color = None
                for run in para.runs:
                    try:
                        if run.font.size:
                            sz = int(run.font.size.pt)
                        if run.font.bold:
                            bold = True
                        color = str(run.font.color.rgb)
                    except:
                        pass
                    break
                print(f"  [{sz}pt bold={bold} color={color}] {line[:120]}")
