"""Generates PP1 Progress Presentation for R26-IT-117."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x23, 0x7E)
BLUE       = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
TEAL       = RGBColor(0x00, 0x89, 0x7B)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN= RGBColor(0xE8, 0xF5, 0xE9)
ORANGE     = RGBColor(0xE6, 0x51, 0x00)
LIGHT_ORANGE=RGBColor(0xFF, 0xF3, 0xE0)
RED        = RGBColor(0xC6, 0x28, 0x28)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x21, 0x21, 0x21)
GREY       = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY   = RGBColor(0x75, 0x75, 0x75)
DARK_GREY  = RGBColor(0x42, 0x42, 0x42)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w if line_w else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    """Dark navy header bar with title."""
    box(slide, 0, 0, W, Inches(1.1), fill=NAVY)
    txt(slide, title, Inches(0.4), Inches(0.12), Inches(10), Inches(0.55),
        size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.65), Inches(10), Inches(0.38),
            size=14, bold=False, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)
    # Project ID right
    txt(slide, "R26-IT-117", Inches(11.8), Inches(0.38), Inches(1.4), Inches(0.4),
        size=11, bold=True, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.RIGHT)


def footer(slide, text="R26-IT-117  |  SLIIT Faculty of Computing  |  Department of Information Technology"):
    box(slide, 0, Inches(7.15), W, Inches(0.35), fill=NAVY)
    txt(slide, text, Inches(0.3), Inches(7.17), Inches(12.7), Inches(0.3),
        size=9, color=WHITE, align=PP_ALIGN.CENTER)


def tag(slide, text, x, y, w=Inches(1.6), h=Inches(0.32), bg=NAVY, fg=WHITE, size=10):
    box(slide, x, y, w, h, fill=bg)
    txt(slide, text, x, y, w, h, size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


def bullet_block(slide, items, x, y, w, h, title=None, title_color=NAVY,
                 bg=GREY, bullet="▸", size=13, title_size=13):
    box(slide, x, y, w, h, fill=bg, line=RGBColor(0xDD, 0xDD, 0xDD), line_w=Pt(0.75))
    ty = y
    if title:
        txt(slide, title, x + Inches(0.12), ty + Inches(0.08),
            w - Inches(0.2), Inches(0.35),
            size=title_size, bold=True, color=title_color)
        ty += Inches(0.38)
    for item in items:
        txt(slide, f"{bullet}  {item}", x + Inches(0.15), ty,
            w - Inches(0.25), Inches(0.32), size=size, color=DARK_GREY)
        ty += Inches(0.3)


def stat_box(slide, value, label, x, y, w=Inches(1.9), h=Inches(1.1),
             val_color=NAVY, bg=LIGHT_BLUE):
    box(slide, x, y, w, h, fill=bg, line=BLUE, line_w=Pt(1))
    txt(slide, value, x, y + Inches(0.08), w, Inches(0.52),
        size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.55), w, Inches(0.48),
        size=10, color=DARK_GREY, align=PP_ALIGN.CENTER, wrap=True)


def stage_card(slide, number, title, desc, x, y, w=Inches(2.9), h=Inches(1.55),
               num_bg=NAVY, card_bg=LIGHT_BLUE):
    box(slide, x, y, w, h, fill=card_bg, line=BLUE, line_w=Pt(0.75))
    # Number badge
    box(slide, x, y, Inches(0.45), h, fill=num_bg)
    txt(slide, number, x, y + Inches(0.52), Inches(0.45), Inches(0.5),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, x + Inches(0.52), y + Inches(0.1),
        w - Inches(0.6), Inches(0.35), size=12, bold=True, color=NAVY)
    txt(slide, desc, x + Inches(0.52), y + Inches(0.42),
        w - Inches(0.6), Inches(1.0), size=10, color=DARK_GREY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title"""
    sl = blank(prs)
    # Full background
    box(sl, 0, 0, W, H, fill=NAVY)
    # White accent strip
    box(sl, 0, Inches(4.8), W, Inches(0.06), fill=WHITE)
    # Light band at bottom
    box(sl, 0, Inches(4.86), W, Inches(2.64), fill=RGBColor(0x0D, 0x16, 0x5E))

    txt(sl, "PROGRESS PRESENTATION | PP1",
        Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.7),
        size=18, bold=True, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

    txt(sl, "AI-Based Intelligent Construction\nPlanning and Management",
        Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.8),
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Project ID: R26-IT-117",
        Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.5),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Shazni M.I.A  |  N.H. Afaaf Hariri  |  Hanfi A.M.M  |  Moorthy T",
        Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.45),
        size=13, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

    txt(sl, "R26-IT-117  |  Sri Lanka Institute of Information Technology  |  May 2026",
        Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.4),
        size=11, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)


def slide_team(prs):
    """Slide 2 — Team"""
    sl = blank(prs)
    header_bar(sl, "Our Team", "R26-IT-117  |  SLIIT Faculty of Computing")

    members = [
        ("Shazni M.I.A",    "IT22131256", "Component 01\nAI Architectural Planning", NAVY),
        ("N.H. Afaaf Hariri","IT22285188","Component 02\nAI Cost Estimation",        BLUE),
        ("Hanfi A.M.M",      "IT22074454", "Component 03\nTimeline & Scheduling",     TEAL),
        ("Moorthy T",        "IT22244734", "Component 04\nPerformance Monitoring",    RGBColor(0x6A,0x1B,0x9A)),
    ]

    xs = [Inches(0.4), Inches(3.65), Inches(6.9), Inches(10.15)]
    w  = Inches(2.9)

    for (name, sid, role, color), x in zip(members, xs):
        # Card
        box(sl, x, Inches(1.4), w, Inches(5.5), fill=WHITE,
            line=color, line_w=Pt(1.5))
        # Top colour band
        box(sl, x, Inches(1.4), w, Inches(0.5), fill=color)
        # Photo placeholder
        box(sl, x + Inches(0.85), Inches(2.05), Inches(1.2), Inches(1.5),
            fill=LIGHT_BLUE, line=color, line_w=Pt(1))
        txt(sl, "Photo", x + Inches(0.85), Inches(2.55), Inches(1.2), Inches(0.4),
            size=10, color=MID_GREY, align=PP_ALIGN.CENTER)
        txt(sl, name, x + Inches(0.1), Inches(3.75), w - Inches(0.2), Inches(0.4),
            size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(sl, sid, x + Inches(0.1), Inches(4.15), w - Inches(0.2), Inches(0.3),
            size=11, color=MID_GREY, align=PP_ALIGN.CENTER)
        txt(sl, role, x + Inches(0.1), Inches(4.55), w - Inches(0.2), Inches(0.6),
            size=10, color=DARK_GREY, align=PP_ALIGN.CENTER)

    # Supervisors
    txt(sl, "Supervisor: Ms. Jenny Krishara  |  Co-Supervisor: Dr. Dinuka Wijendra",
        Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
        size=11, color=MID_GREY, align=PP_ALIGN.CENTER)
    footer(sl)


def slide_agenda(prs):
    """Slide 3 — Agenda"""
    sl = blank(prs)
    header_bar(sl, "Agenda", "Progress Presentation PP1  |  R26-IT-117")

    items = [
        ("01", "Project Overview & Architecture",    "Brief recap of the overall AI pipeline"),
        ("02", "Component 01 — AI Architectural Planning",  "Stage 1–4 implementation, models, metrics, demo"),
        ("03", "Component 02 — AI Cost Estimation",  "Progress update"),
        ("04", "Component 03 — Timeline & Scheduling","Progress update"),
        ("05", "Component 04 — Performance Monitoring","Progress update"),
        ("06", "Overall Progress & Next Steps",       "Milestones, challenges, roadmap"),
    ]

    colors = [NAVY, RGBColor(0x1B,0x5E,0x20), BLUE, TEAL,
              RGBColor(0x6A,0x1B,0x9A), RGBColor(0xBF,0x36,0x0C)]

    for i, ((num, title, desc), color) in enumerate(zip(items, colors)):
        row = i // 2
        col = i % 2
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.35) + row * Inches(1.85)
        w = Inches(6.1)
        h = Inches(1.6)
        box(sl, x, y, w, h, fill=WHITE, line=color, line_w=Pt(1.5))
        box(sl, x, y, Inches(0.6), h, fill=color)
        txt(sl, num, x, y + Inches(0.55), Inches(0.6), Inches(0.5),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, title, x + Inches(0.7), y + Inches(0.2), w - Inches(0.8), Inches(0.45),
            size=13, bold=True, color=color)
        txt(sl, desc, x + Inches(0.7), y + Inches(0.65), w - Inches(0.8), Inches(0.8),
            size=11, color=DARK_GREY)
    footer(sl)


def slide_system_overview(prs):
    """Slide 4 — System Architecture Overview"""
    sl = blank(prs)
    header_bar(sl, "System Architecture Overview", "AI-Based Intelligent Construction Planning and Management")

    # Input box
    box(sl, Inches(0.3), Inches(1.35), Inches(1.9), Inches(5.5),
        fill=LIGHT_BLUE, line=BLUE, line_w=Pt(1))
    txt(sl, "INPUT", Inches(0.3), Inches(1.35), Inches(1.9), Inches(0.4),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sl, Inches(0.3), Inches(1.35), Inches(1.9), Inches(0.4), fill=BLUE)
    txt(sl, "INPUT", Inches(0.3), Inches(1.35), Inches(1.9), Inches(0.4),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for item, yoff in [("Cadastral Plan\n(PDF/JPG)", 0.55), ("User\nRequirements", 1.4), ("Progress\nUpdates", 2.3)]:
        txt(sl, item, Inches(0.35), Inches(1.35) + Inches(yoff), Inches(1.8), Inches(0.7),
            size=10, color=NAVY, align=PP_ALIGN.CENTER)

    # 4 components
    comp_data = [
        ("C01", "AI Architectural\nPlanning", "OCR → NER → Buildable Zone\n→ LLM-RAG → SVG/PDF", NAVY, Inches(1.5)),
        ("C02", "AI Cost\nEstimation", "BOQ Generator → ICTAD\nRates → Ensemble ML", BLUE, Inches(3.1)),
        ("C03", "Timeline &\nScheduling", "LSTM + Genetic Algorithm\n→ Gantt Chart", TEAL, Inches(4.7)),
        ("C04", "Performance\nMonitoring", "SPI Monitor → XGBoost\n→ RAG + LLM", RGBColor(0x6A,0x1B,0x9A), Inches(6.3)),
    ]

    for cid, name, tech, color, y in comp_data:
        box(sl, Inches(2.5), y, Inches(7.5), Inches(1.35), fill=WHITE,
            line=color, line_w=Pt(1.5))
        box(sl, Inches(2.5), y, Inches(0.7), Inches(1.35), fill=color)
        txt(sl, cid, Inches(2.5), y + Inches(0.42), Inches(0.7), Inches(0.5),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, name, Inches(3.3), y + Inches(0.08), Inches(2.2), Inches(0.6),
            size=12, bold=True, color=color)
        txt(sl, tech, Inches(3.3), y + Inches(0.65), Inches(6.5), Inches(0.6),
            size=10, color=DARK_GREY)

    # Output box
    box(sl, Inches(10.3), Inches(1.35), Inches(2.7), Inches(5.5),
        fill=LIGHT_GREEN, line=GREEN, line_w=Pt(1))
    box(sl, Inches(10.3), Inches(1.35), Inches(2.7), Inches(0.4), fill=GREEN)
    txt(sl, "OUTPUT", Inches(10.3), Inches(1.35), Inches(2.7), Inches(0.4),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for item, yoff in [("Floor Plans\n(SVG + PDF)", 0.55), ("Cost Estimate\n(JSON + PDF)", 1.4),
                        ("Project\nSchedule", 2.3), ("Delay\nAlerts", 3.15)]:
        txt(sl, item, Inches(10.35), Inches(1.35) + Inches(yoff), Inches(2.6), Inches(0.7),
            size=10, color=RGBColor(0x1B,0x5E,0x20), align=PP_ALIGN.CENTER)

    footer(sl)


def slide_progress_summary(prs):
    """Slide 5 — Overall Progress Summary"""
    sl = blank(prs)
    header_bar(sl, "Overall Progress Summary", "Sprint Status — May 2026")

    rows = [
        ("Component 01", "AI Architectural Planning",  "Shazni M.I.A",    "Complete",     GREEN,  "100%"),
        ("Component 02", "AI Cost Estimation",         "Afaaf Hariri",     "In Progress",  BLUE,   "40%"),
        ("Component 03", "Timeline & Scheduling",      "Hanfi A.M.M",      "In Progress",  TEAL,   "35%"),
        ("Component 04", "Performance Monitoring",     "Moorthy T",        "In Progress",  RGBColor(0x6A,0x1B,0x9A), "30%"),
    ]

    headers = ["Component", "Title", "Owner", "Status", "Progress"]
    col_ws  = [Inches(1.6), Inches(3.8), Inches(2.0), Inches(1.6), Inches(3.8)]
    col_xs  = [Inches(0.3)]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w)

    # Header row
    y = Inches(1.3)
    for hdr, cx, cw in zip(headers, col_xs, col_ws):
        box(sl, cx, y, cw, Inches(0.42), fill=NAVY)
        txt(sl, hdr, cx + Inches(0.08), y + Inches(0.05), cw - Inches(0.1), Inches(0.35),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (cid, title, owner, status, color, pct) in enumerate(rows):
        y = Inches(1.72) + i * Inches(1.22)
        row_bg = WHITE if i % 2 == 0 else GREY
        for cx, cw in zip(col_xs, col_ws):
            box(sl, cx, y, cw, Inches(1.12), fill=row_bg, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))

        txt(sl, cid, col_xs[0] + Inches(0.08), y + Inches(0.35), col_ws[0], Inches(0.4),
            size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(sl, title, col_xs[1] + Inches(0.1), y + Inches(0.35), col_ws[1], Inches(0.4),
            size=12, color=BLACK)
        txt(sl, owner, col_xs[2] + Inches(0.08), y + Inches(0.35), col_ws[2], Inches(0.4),
            size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)

        # Status badge
        badge_bg = LIGHT_GREEN if status == "Complete" else LIGHT_ORANGE
        badge_fg = GREEN if status == "Complete" else ORANGE
        box(sl, col_xs[3] + Inches(0.2), y + Inches(0.28), Inches(1.2), Inches(0.35), fill=badge_bg)
        txt(sl, status, col_xs[3] + Inches(0.2), y + Inches(0.28), Inches(1.2), Inches(0.35),
            size=10, bold=True, color=badge_fg, align=PP_ALIGN.CENTER)

        # Progress bar
        bar_x = col_xs[4] + Inches(0.15)
        bar_y = y + Inches(0.45)
        bar_w = Inches(3.4)
        bar_h = Inches(0.22)
        box(sl, bar_x, bar_y, bar_w, bar_h, fill=RGBColor(0xE0,0xE0,0xE0))
        fill_w = bar_w * int(pct.replace("%","")) // 100
        box(sl, bar_x, bar_y, fill_w, bar_h, fill=color)
        txt(sl, pct, bar_x + bar_w + Inches(0.08), bar_y - Inches(0.02), Inches(0.5), bar_h,
            size=11, bold=True, color=color)

    footer(sl)


def slide_c01_intro(prs):
    """Slide 6 — C01 Component Introduction"""
    sl = blank(prs)
    box(sl, 0, 0, W, H, fill=NAVY)
    box(sl, 0, Inches(5.5), W, Inches(2.0), fill=RGBColor(0x0D,0x16,0x5E))

    tag(sl, "COMPONENT 01", Inches(0.5), Inches(0.5),
        w=Inches(2.2), h=Inches(0.42), bg=WHITE, fg=NAVY, size=12)

    txt(sl, "AI Architectural Planning Model",
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.2),
        size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    txt(sl, "First cadastral-to-floor-plan pipeline for Sri Lanka using LLM-RAG architecture",
        Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.6),
        size=16, color=RGBColor(0xBB,0xDE,0xFB), align=PP_ALIGN.LEFT)

    details = [
        "FastAPI 4-stage pipeline: Extract → Buildable Zone → Generate → Render",
        "Gemini 2.5 Flash + RAG (131 Sri Lankan floor plan entries in Supabase pgvector)",
        "spaCy NER fine-tuned on 69 real Sri Lankan cadastral plans — F-score: 0.857",
        "Output: 3 SVG/PDF floor plan alternatives + Building Schema JSON for C02",
    ]
    for i, d in enumerate(details):
        txt(sl, f"✓  {d}", Inches(0.5), Inches(3.15) + i * Inches(0.48),
            Inches(12.3), Inches(0.42), size=13, color=WHITE)

    txt(sl, "Shazni M.I.A  |  IT22131256  |  BSc (Hons) Information Technology",
        Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
        size=12, color=RGBColor(0x90,0xCA,0xF9), align=PP_ALIGN.LEFT)
    txt(sl, "STATUS: COMPLETE", Inches(10.2), Inches(5.8), Inches(2.6), Inches(0.4),
        size=13, bold=True, color=RGBColor(0x69,0xF0,0xAE), align=PP_ALIGN.RIGHT)


def slide_c01_pipeline(prs):
    """Slide 7 — C01 4-Stage Pipeline"""
    sl = blank(prs)
    header_bar(sl, "C01 — 4-Stage Pipeline Overview",
               "AI Architectural Planning Model  |  Shazni M.I.A  |  IT22131256")

    stages = [
        ("01", "Stage 1\nExtraction",
         "• EasyOCR (CRAFT+CRNN)\n• spaCy NER (F=0.857)\n• OpenCV boundary detection\n• pyproj SLD99→WGS84\n• CNN cadastral classifier",
         "Site Schema JSON", NAVY),
        ("02", "Stage 2\nBuildable Zone",
         "• NBC Sri Lanka setbacks\n• BCR by district & coastal flag\n• Shapely polygon buffering\n• Orientation solver\n• Coordinate scale calibration",
         "Buildable Polygon + Footprint", BLUE),
        ("03", "Stage 3\nFloor Plan Gen",
         "• RAG: Supabase pgvector\n  131 SL floor plan entries\n• Gemini 2.5 Flash × 3 temps\n• Strip-packing layout solver\n• Geometric validator + scorer",
         "3 Floor Plan Alternatives", TEAL),
        ("04", "Stage 4\nRender",
         "• SVGRenderer (svgwrite)\n• PDFRenderer (ReportLab)\n• SchemaSerialiser → C02\n• Building Schema JSON\n• jsonschema validation",
         "SVG + PDF + Schema JSON", GREEN),
    ]

    for i, (num, title, desc, output, color) in enumerate(stages):
        x = Inches(0.28) + i * Inches(3.26)
        w = Inches(3.0)
        # Card
        box(sl, x, Inches(1.3), w, Inches(5.5), fill=WHITE, line=color, line_w=Pt(1.5))
        # Top badge
        box(sl, x, Inches(1.3), w, Inches(0.75), fill=color)
        txt(sl, num, x + Inches(0.1), Inches(1.32), Inches(0.5), Inches(0.45),
            size=18, bold=True, color=WHITE)
        txt(sl, title, x + Inches(0.55), Inches(1.35), w - Inches(0.6), Inches(0.65),
            size=12, bold=True, color=WHITE)
        # Description
        txt(sl, desc, x + Inches(0.15), Inches(2.15), w - Inches(0.25), Inches(3.2),
            size=11, color=DARK_GREY)
        # Output tag
        box(sl, x + Inches(0.1), Inches(6.35), w - Inches(0.2), Inches(0.35), fill=color)
        txt(sl, f"→ {output}", x + Inches(0.1), Inches(6.35), w - Inches(0.2), Inches(0.35),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow between stages
        if i < 3:
            txt(sl, "▶", x + w + Inches(0.05), Inches(3.8), Inches(0.22), Inches(0.4),
                size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    footer(sl)


def slide_c01_stage1(prs):
    """Slide 8 — Stage 1 Detail"""
    sl = blank(prs)
    header_bar(sl, "C01 — Stage 1: Cadastral Plan Extraction",
               "POST /api/v1/extract  |  Synchronous")

    # Left column
    bullet_block(sl, [
        "PDF/image upload via FastAPI multipart",
        "PyMuPDF renders PDF page to image",
        "EasyOCR (CRAFT + CRNN) reads all text",
        "OpenCV (Canny + contours) detects boundary polygon",
        "CNN classifier rejects non-cadastral uploads",
    ], Inches(0.3), Inches(1.3), Inches(4.1), Inches(2.2),
       title="OCR & Boundary Detection", bg=LIGHT_BLUE, title_color=NAVY)

    bullet_block(sl, [
        "plan_number, district, area_sqm, scale",
        "surveyor, road_access, lot_number",
        "coordinate_n, coordinate_e (SLD99 range filter)",
        "licence_number, is_coastal flag",
        "Fine-tuned on 69 real Sri Lankan cadastral plans",
        "F-score: 0.857  |  Precision: 0.941  |  Recall: 0.787",
    ], Inches(0.3), Inches(3.65), Inches(4.1), Inches(2.85),
       title="spaCy NER — Extracted Fields", bg=LIGHT_BLUE, title_color=NAVY)

    # Right column
    bullet_block(sl, [
        "SLD99 (EPSG:5234) → WGS84 (EPSG:4326)",
        "pyproj handles mathematical transformation",
        "Range filter: N=100K–600K, E=200K–700K",
        "Prevents phone numbers matching coordinates",
        "Example: E=617050 N=539050 → lat=9.35 lon=80.35",
    ], Inches(4.65), Inches(1.3), Inches(4.2), Inches(2.2),
       title="Coordinate Conversion (SLD99 → GPS)", bg=LIGHT_GREEN, title_color=GREEN)

    bullet_block(sl, [
        "5 classes: surveyor_plan, municipal_plan,",
        "  uda_plan, low_quality, non_cadastral",
        "ResNet-based CNN — fine-tuning in progress",
        "Collecting 100+ real plans from licensed surveyors",
        "Expected accuracy: 90%+",
    ], Inches(4.65), Inches(3.65), Inches(4.2), Inches(2.0),
       title="CNN Cadastral Classifier", bg=LIGHT_GREEN, title_color=GREEN)

    # Output schema
    box(sl, Inches(9.1), Inches(1.3), Inches(4.0), Inches(5.2),
        fill=GREY, line=NAVY, line_w=Pt(1))
    txt(sl, "Site Schema JSON Output", Inches(9.2), Inches(1.35), Inches(3.8), Inches(0.38),
        size=12, bold=True, color=NAVY)
    schema_txt = (
        '{\n'
        '  "plan_id": "PLAN-2720-AM-...",\n'
        '  "district": "Ampara",\n'
        '  "area_sqm": 841.0,\n'
        '  "is_coastal": true,\n'
        '  "road_access": "lane",\n'
        '  "coordinate_n": 538975.0,\n'
        '  "coordinate_e": 617050.0,\n'
        '  "gps_lat": 9.352,\n'
        '  "gps_lon": 80.358,\n'
        '  "boundary_polygon": [...]\n'
        '}'
    )
    txt(sl, schema_txt, Inches(9.2), Inches(1.78), Inches(3.8), Inches(4.5),
        size=10, color=DARK_GREY)
    footer(sl)


def slide_c01_stage2(prs):
    """Slide 9 — Stage 2 Detail"""
    sl = blank(prs)
    header_bar(sl, "C01 — Stage 2: Buildable Zone Calculation",
               "POST /api/v1/buildable-zone  |  Synchronous  |  NBC Sri Lanka")

    # NBC rules
    bullet_block(sl, [
        "National Road / Main Road → Front: 3.0 m",
        "Provincial Road → Front: 2.0 m",
        "Local Road → Front: 1.5 m",
        "Lane / Private Road → Front: 1.0 m",
        "Rear setback: 1.5 m minimum",
        "Side setback: 1.0 m minimum",
    ], Inches(0.3), Inches(1.3), Inches(4.0), Inches(2.6),
       title="NBC Sri Lanka — Setback Rules", bg=LIGHT_BLUE, title_color=NAVY)

    bullet_block(sl, [
        "Colombo (urban): BCR = 0.60",
        "Other districts (standard): BCR = 0.50",
        "Coastal districts (Ampara, Galle, etc.): BCR = 0.40",
        "Coastal → Pile foundation (depth: 2.5 m)",
        "Large plot >500 sqm → Raft foundation (1.8 m)",
        "Standard plot → Strip foundation (1.5 m)",
    ], Inches(0.3), Inches(4.1), Inches(4.0), Inches(2.75),
       title="BCR & Foundation Rules", bg=LIGHT_BLUE, title_color=NAVY)

    # Shapely
    bullet_block(sl, [
        "Normalised polygon coords (0–1) from Stage 1",
        "scale_factor = sqrt(real_area / norm_area)",
        "Converts setback metres → normalised units",
        "Shapely .buffer(-setback_norm, join_style=2)",
        "All output areas converted back to real sqm",
        "Fallback: halve setback if polygon becomes empty",
    ], Inches(4.55), Inches(1.3), Inches(4.2), Inches(2.6),
       title="Shapely Polygon Calculation", bg=LIGHT_GREEN, title_color=GREEN)

    # Example output
    box(sl, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.5),
        fill=GREY, line=NAVY, line_w=Pt(1))
    txt(sl, "Example Output — Plan 2720 (Ampara)",
        Inches(9.1), Inches(1.35), Inches(3.8), Inches(0.38),
        size=11, bold=True, color=NAVY)

    stats = [
        ("Plot Area", "841.0 sqm"),
        ("BCR Applied", "0.40 (coastal)"),
        ("Max Footprint", "336.4 sqm"),
        ("Max Built (2F)", "672.8 sqm"),
        ("Rec. Floors", "3"),
        ("Entrance Side", "East (lane)"),
        ("Foundation", "Pile (coastal)"),
    ]
    for i, (k, v) in enumerate(stats):
        y = Inches(1.85) + i * Inches(0.6)
        box(sl, Inches(9.1), y, Inches(1.9), Inches(0.45), fill=LIGHT_BLUE)
        txt(sl, k, Inches(9.15), y + Inches(0.05), Inches(1.8), Inches(0.38),
            size=10, color=NAVY)
        box(sl, Inches(9.1) + Inches(1.9), y, Inches(1.9), Inches(0.45), fill=WHITE,
            line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
        txt(sl, v, Inches(9.15) + Inches(1.9), y + Inches(0.05), Inches(1.8), Inches(0.38),
            size=10, bold=True, color=DARK_GREY)

    footer(sl)


def slide_c01_stage3(prs):
    """Slide 10 — Stage 3 Detail"""
    sl = blank(prs)
    header_bar(sl, "C01 — Stage 3: AI Floor Plan Generation",
               "POST /api/v1/generate-floor-plan  |  Async via Celery + Redis")

    # RAG
    bullet_block(sl, [
        "Vector DB: Supabase pgvector (cloud PostgreSQL)",
        "131 entries: 10 SL coastal + 121 CVC FPP dataset",
        "Embedder: all-MiniLM-L6-v2 (384-dim vectors)",
        "Retrieves 5 most similar designs — cosine similarity",
        "Gives Gemini real Sri Lankan context",
        "Without RAG: generic Western-style plans generated",
    ], Inches(0.3), Inches(1.3), Inches(4.2), Inches(2.9),
       title="RAG — Supabase pgvector", bg=LIGHT_BLUE, title_color=NAVY)

    # Gemini
    bullet_block(sl, [
        "Model: Gemini 2.5 Flash (google-genai SDK)",
        "Temperature 0.4 → Conservative (strict NBC)",
        "Temperature 0.7 → Balanced (rules + creativity)",
        "Temperature 1.0 → Creative (innovative design)",
        "3 independent API calls — 3 real alternatives",
        "response_mime_type: application/json",
    ], Inches(0.3), Inches(4.35), Inches(4.2), Inches(2.75),
       title="Gemini 2.5 Flash — 3 Temperatures", bg=LIGHT_BLUE, title_color=NAVY)

    # Layout solver + validator
    bullet_block(sl, [
        "Strip-packing solver: repositions rooms to eliminate overlaps",
        "Places rooms within buildable zone [x_min, y_min]–[x_max, y_max]",
        "Validator: checks NBC minimum room sizes",
        "  living_room ≥16, master_bedroom ≥12, bedroom ≥9 sqm",
        "Scorer: 4 dimensions — space utilisation, natural light,",
        "  adjacency quality, ventilation potential",
    ], Inches(4.65), Inches(1.3), Inches(8.35), Inches(2.9),
       title="Layout Solver + Validator + Scorer", bg=LIGHT_GREEN, title_color=GREEN)

    # Results
    stats_data = [
        ("Task Queue", "Celery 5.4 + Redis"),
        ("Generation Time", "~145 seconds (3 calls)"),
        ("RAG Retrieval", "5 passages per query"),
        ("Validity (841 sqm)", "3/3 is_valid = true ✓"),
        ("Quality Score", "0.88 / 1.0 (balanced)"),
    ]
    box(sl, Inches(4.65), Inches(4.35), Inches(8.35), Inches(2.75),
        fill=GREY, line=GREEN, line_w=Pt(1))
    txt(sl, "Stage 3 Results — Plan 2720 (841 sqm, 2 floors)",
        Inches(4.75), Inches(4.4), Inches(8.1), Inches(0.38),
        size=11, bold=True, color=GREEN)
    for i, (k, v) in enumerate(stats_data):
        y = Inches(4.9) + i * Inches(0.42)
        txt(sl, f"▸  {k}:", Inches(4.8), y, Inches(2.5), Inches(0.38), size=11, color=DARK_GREY)
        txt(sl, v, Inches(7.2), y, Inches(5.6), Inches(0.38), size=11, bold=True, color=NAVY)

    footer(sl)


def slide_c01_stage4(prs):
    """Slide 11 — Stage 4 Detail"""
    sl = blank(prs)
    header_bar(sl, "C01 — Stage 4: Rendering & Schema Output",
               "POST /api/v1/render  |  Synchronous  |  SVG + PDF + Building Schema JSON")

    bullet_block(sl, [
        "svgwrite library — Python SVG generation",
        "Thick walls (4px) + distinct room colours",
        "Door arcs per window orientation",
        "Furniture symbols: bed, sofa, kitchen counter, toilet",
        "Dimension lines (width + height in metres)",
        "Window indicators (cyan bars on walls)",
        "North arrow + scale bar + legend + title block",
        "Multi-floor: Ground Floor + Floor 2 side by side",
    ], Inches(0.3), Inches(1.3), Inches(4.2), Inches(4.7),
       title="SVG Renderer (svgwrite)", bg=LIGHT_BLUE, title_color=NAVY)

    bullet_block(sl, [
        "ReportLab: 6-page professional PDF report",
        "Cover page + floor plans + site summary",
        "Building schema summary + quality scores",
    ], Inches(4.65), Inches(1.3), Inches(4.0), Inches(1.4),
       title="PDF Renderer (ReportLab)", bg=LIGHT_GREEN, title_color=GREEN)

    bullet_block(sl, [
        "Flat JSON contract for Component 02",
        "footprint_sqm, perimeter, column_count",
        "excavation_depth, openings_sqm, internal_wall_length",
        "finish_grade, roof_type, foundation_type",
        "rooms (count dict), bathroom_count, room_count",
        "base_rate_date, target_date, district, is_coastal",
        "Validated by jsonschema before POST to C02",
        "Fire-and-forget POST to C02_BASE_URL/estimate",
    ], Inches(4.65), Inches(2.85), Inches(4.0), Inches(3.15),
       title="Building Schema JSON → Component 02", bg=LIGHT_GREEN, title_color=GREEN)

    # Sample schema
    box(sl, Inches(8.9), Inches(1.3), Inches(4.1), Inches(5.5),
        fill=GREY, line=NAVY, line_w=Pt(1))
    txt(sl, "Building Schema — Plan 2720",
        Inches(9.0), Inches(1.35), Inches(3.9), Inches(0.38),
        size=11, bold=True, color=NAVY)
    schema = (
        '{\n'
        '  "footprint_sqm": 336.4,\n'
        '  "perimeter": 74.4,\n'
        '  "floors": 2,\n'
        '  "excavation_depth": 2.5,\n'
        '  "column_count": 38,\n'
        '  "openings_sqm": 35.5,\n'
        '  "internal_wall_length": 88.2,\n'
        '  "finish_grade": "mid",\n'
        '  "roof_type": "flat",\n'
        '  "is_coastal": true,\n'
        '  "foundation": "pile",\n'
        '  "rooms": {\n'
        '    "master_bedroom": 1,\n'
        '    "bedroom": 3, "bathroom": 2\n'
        '  },\n'
        '  "room_count": 8\n'
        '}'
    )
    txt(sl, schema, Inches(9.0), Inches(1.78), Inches(3.9), Inches(4.8),
        size=9, color=DARK_GREY)
    footer(sl)


def slide_c01_metrics(prs):
    """Slide 12 — C01 Accuracy & Metrics"""
    sl = blank(prs)
    header_bar(sl, "C01 — Accuracy & Performance Metrics",
               "AI Architectural Planning Model  |  Evaluation Results")

    metrics = [
        ("0.857", "NER F-Score\n(fine-tuned model)", NAVY),
        ("0.941", "NER Precision\n(P score)", BLUE),
        ("0.787", "NER Recall\n(R score)", TEAL),
        ("100%", "Floor Plan\nValidity Rate", GREEN),
        ("0.88", "Quality Score\n(841 sqm plan)", RGBColor(0xE6,0x51,0x00)),
        ("31/31", "Unit Tests\nPassing", RGBColor(0x6A,0x1B,0x9A)),
    ]

    xs = [Inches(0.3), Inches(2.5), Inches(4.7), Inches(6.9), Inches(9.1), Inches(11.3)]
    for (val, label, color), x in zip(metrics, xs):
        stat_box(sl, val, label, x, Inches(1.4),
                 w=Inches(1.85), h=Inches(1.5), val_color=color,
                 bg=LIGHT_BLUE if color in [NAVY, BLUE, TEAL] else LIGHT_GREEN)

    # NER detail
    bullet_block(sl, [
        "Base model: en_core_web_sm (spaCy)",
        "Fine-tuned on 69 real Sri Lankan cadastral plans",
        "Augmented with 41 synthetic coastal examples",
        "Custom labels: DISTRICT, PLAN_NO, SURVEYOR, LOT_NO, LICENCE",
        "Fallback to base model if fine-tuned model not found",
        "SLD99 range filter prevents phone number false matches",
    ], Inches(0.3), Inches(3.15), Inches(4.1), Inches(3.55),
       title="NER Model — Training Details", bg=LIGHT_BLUE, title_color=NAVY)

    bullet_block(sl, [
        "RAG: 131 entries — 5 passages retrieved per query",
        "Gemini 2.5 Flash: 3 calls at temp 0.4 / 0.7 / 1.0",
        "Layout solver: strip-packing, eliminates all overlaps",
        "Validator: NBC minimum room size enforcement",
        "Scorer: 4 dimensions (space, light, adjacency, ventilation)",
        "841 sqm plan: all 3 alternatives is_valid=true, score=0.88",
    ], Inches(4.55), Inches(3.15), Inches(4.3), Inches(3.55),
       title="Floor Plan Generation — Key Results", bg=LIGHT_GREEN, title_color=GREEN)

    bullet_block(sl, [
        "FastAPI + Uvicorn — port 8001",
        "Celery 5.4 + Redis — async task queue",
        "Generation time: ~145 seconds (3 Gemini calls)",
        "End-to-end pipeline: ~160 seconds total",
        "Lazy ML imports — no startup delay",
        "Fire-and-forget C02 POST (never blocks render)",
    ], Inches(9.1), Inches(3.15), Inches(4.0), Inches(3.55),
       title="Infrastructure & Performance", bg=GREY, title_color=DARK_GREY)

    footer(sl)


def slide_c01_demo(prs):
    """Slide 13 — C01 Demo Output"""
    sl = blank(prs)
    header_bar(sl, "C01 — Demo: End-to-End Pipeline Output",
               "Plan 2720  |  Ampara District  |  841 sqm  |  2 Floors  |  Balanced Layout")

    # Pipeline flow
    steps = [
        ("1", "Upload\nPlan 2720", NAVY),
        ("2", "Extract\nSite Schema", BLUE),
        ("3", "Buildable\nZone 336 sqm", TEAL),
        ("4", "Generate\n3 Alternatives", GREEN),
        ("5", "Render\nSVG + PDF", ORANGE),
        ("6", "Schema\n→ C02", RGBColor(0x6A,0x1B,0x9A)),
    ]
    for i, (num, label, color) in enumerate(steps):
        x = Inches(0.25) + i * Inches(2.17)
        box(sl, x, Inches(1.25), Inches(1.9), Inches(0.9), fill=color)
        txt(sl, num, x + Inches(0.05), Inches(1.28), Inches(0.4), Inches(0.4),
            size=16, bold=True, color=WHITE)
        txt(sl, label, x + Inches(0.1), Inches(1.25), Inches(1.8), Inches(0.9),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 5:
            txt(sl, "▶", x + Inches(1.9), Inches(1.52), Inches(0.27), Inches(0.35),
                size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Results boxes
    results = [
        ("Stage 1 Output", [
            "plan_id: PLAN-2720-AM-20260510",
            "district: Ampara",
            "area_sqm: 841.0",
            "is_coastal: true",
            "road_access: lane",
            "coordinate_n: 538975.0",
        ], NAVY, LIGHT_BLUE),
        ("Stage 2 Output", [
            "max_footprint_sqm: 336.4",
            "coverage_achieved: 0.40",
            "recommended_floors: 3",
            "entrance_side: east",
            "foundation: pile",
            "setbacks: F1.0 R1.5 S1.0",
        ], BLUE, LIGHT_BLUE),
        ("Stage 3 Output", [
            "3 alternatives generated",
            "conservative: score 0.88",
            "balanced: score 0.88",
            "creative: score 0.78",
            "all is_valid = true",
            "RAG: 5 passages from Supabase",
        ], GREEN, LIGHT_GREEN),
        ("Stage 4 Output", [
            "SVG: 2-floor plan rendered",
            "rooms: 10 (5 per floor)",
            "column_count: 38",
            "perimeter: 74.4 m",
            "Building Schema → C02",
            "jsonschema validation: PASS",
        ], TEAL, LIGHT_GREEN),
    ]

    for i, (title, items, color, bg) in enumerate(results):
        x = Inches(0.3) + i * Inches(3.25)
        bullet_block(sl, items, x, Inches(2.35), Inches(3.0), Inches(4.7),
                     title=title, bg=bg, title_color=color, size=10)

    footer(sl)


def slide_c01_design_excellence(prs):
    """Slide 14 — C01 Design Excellence"""
    sl = blank(prs)
    header_bar(sl, "C01 — Design Excellence & Key Achievements",
               "AI Architectural Planning Model  |  Component 01")

    achievements = [
        ("Sri Lanka-First", "First cadastral-to-floor-plan pipeline for Sri Lanka. No existing system converts land survey documents directly into optimized floor plans.", NAVY),
        ("NBC Compliance", "Full NBC Sri Lanka enforcement: setbacks, BCR per district, coastal pile foundation rules. All outputs comply with national building code.", BLUE),
        ("Real Dataset", "NER trained on 69 real cadastral plans from Ampara district surveyors. RAG store has 131 entries from real Sri Lankan house designs.", TEAL),
        ("End-to-End API", "FastAPI REST API with 7 endpoints. Async generation via Celery + Redis. Swagger docs at /docs. 31/31 unit tests passing.", GREEN),
        ("Flat C02 Contract", "Building Schema JSON matches exact flat format expected by Component 02. Validated by jsonschema. POST is fire-and-forget (non-blocking).", ORANGE),
        ("google-genai SDK", "Migrated from deprecated google.generativeai to google.genai. Supports Gemini 2.5 Flash with JSON mode and AFC function calling.", RGBColor(0x6A,0x1B,0x9A)),
    ]

    for i, (title, desc, color) in enumerate(achievements):
        row = i // 3
        col = i % 3
        x = Inches(0.3) + col * Inches(4.33)
        y = Inches(1.3) + row * Inches(2.55)
        w, h = Inches(4.05), Inches(2.35)
        box(sl, x, y, w, h, fill=WHITE, line=color, line_w=Pt(1.5))
        box(sl, x, y, w, Inches(0.42), fill=color)
        txt(sl, title, x + Inches(0.1), y + Inches(0.05), w - Inches(0.15), Inches(0.35),
            size=13, bold=True, color=WHITE)
        txt(sl, desc, x + Inches(0.12), y + Inches(0.52), w - Inches(0.2), Inches(1.7),
            size=11, color=DARK_GREY, wrap=True)

    footer(sl)


def slide_c01_challenges(prs):
    """Slide 15 — C01 Challenges & Solutions"""
    sl = blank(prs)
    header_bar(sl, "C01 — Challenges & Solutions",
               "AI Architectural Planning Model  |  Component 01")

    challenges = [
        (
            "Normalised Polygon Scale Issue",
            "Stage 1 outputs polygon coords normalised (0–1). Stage 2 was treating them as real metres, causing wrong setback distances.",
            "Added scale_factor = sqrt(real_area / norm_area) to convert setback metres → normalised units. All areas now correctly in sqm.",
            NAVY,
        ),
        (
            "Phone Number → Coordinate False Match",
            "NER regex matched '077' from phone number '077 6550260' as the E coordinate instead of real SLD99 easting '617050 E'.",
            "Changed from regex search() to finditer() with SLD99 range validation: N=100K–600K, E=200K–700K. Phone numbers rejected.",
            BLUE,
        ),
        (
            "Supabase RPC Type Mismatch",
            "match_sl_plans RPC function returned bigint but the id column was changed to uuid, causing RAG to fall back to 3 hardcoded norms.",
            "Dropped and recreated the RPC function with correct uuid return type. RAG now retrieves 5 real passages from 131 Supabase entries.",
            TEAL,
        ),
        (
            "Building Schema JSON Contract",
            "building_schema.json validated old nested structure (site{}, buildable_zone{}, building{}). SchemaSerialiser was rewritten for flat C02 format.",
            "Updated building_schema.json to validate the flat C02 contract. All required fields match exactly what Component 02 expects.",
            GREEN,
        ),
    ]

    for i, (title, problem, solution, color) in enumerate(challenges):
        row = i // 2
        col = i % 2
        x = Inches(0.3) + col * Inches(6.5)
        y = Inches(1.3) + row * Inches(2.85)
        w, h = Inches(6.2), Inches(2.65)
        box(sl, x, y, w, h, fill=WHITE, line=color, line_w=Pt(1.5))
        box(sl, x, y, w, Inches(0.4), fill=color)
        txt(sl, title, x + Inches(0.1), y + Inches(0.05), w - Inches(0.15), Inches(0.35),
            size=12, bold=True, color=WHITE)
        txt(sl, "Problem:", x + Inches(0.1), y + Inches(0.5), Inches(0.8), Inches(0.28),
            size=10, bold=True, color=RED)
        txt(sl, problem, x + Inches(0.9), y + Inches(0.5), w - Inches(1.0), Inches(0.7),
            size=10, color=DARK_GREY, wrap=True)
        txt(sl, "Solution:", x + Inches(0.1), y + Inches(1.3), Inches(0.85), Inches(0.28),
            size=10, bold=True, color=GREEN)
        txt(sl, solution, x + Inches(0.95), y + Inches(1.3), w - Inches(1.05), Inches(1.2),
            size=10, color=DARK_GREY, wrap=True)

    footer(sl)


def slide_next_steps(prs):
    """Slide 16 — Next Steps"""
    sl = blank(prs)
    header_bar(sl, "Next Steps & Roadmap", "R26-IT-117  |  Sprint Planning")

    items = [
        ("C01 — Shazni", [
            "CNN classifier fine-tuning (100+ real cadastral plans)",
            "Expand NER training data to all 25 districts",
            "Migrate to directional Shapely buffers (Sprint 5)",
            "Solar path simulation per district latitude (Sprint 6)",
            "Frontend integration with Component 04",
        ], NAVY, LIGHT_BLUE),
        ("C02 — Afaaf Hariri", [
            "Complete ICTAD rate database (district-level)",
            "Train Ensemble ML (XGBoost 65% + MLP 35%)",
            "Achieve MAPE < 12% on 60+ historical projects",
            "BOQ Generator + Risk Adjuster implementation",
            "Consume Building Schema JSON from C01",
        ], BLUE, LIGHT_BLUE),
        ("C03 — Hanfi", [
            "Collect historical construction project data",
            "Train LSTM model for phase duration prediction",
            "Genetic Algorithm for resource optimization",
            "Dynamic Gantt chart generation",
            "Integrate with C02 cost output",
        ], TEAL, LIGHT_GREEN),
        ("C04 — Moorthy", [
            "SPI-based real-time progress monitoring",
            "Train XGBoost delay prediction model (F1 ≥ 0.80)",
            "RAG pipeline: FAISS + Sentence Transformers",
            "LLM corrective recommendations (Gemini)",
            "Dashboard with delay alerts",
        ], RGBColor(0x6A,0x1B,0x9A), LIGHT_GREEN),
    ]

    for i, (title, bullets, color, bg) in enumerate(items):
        x = Inches(0.3) + i * Inches(3.25)
        bullet_block(sl, bullets, x, Inches(1.3), Inches(3.0), Inches(5.6),
                     title=title, bg=bg, title_color=color, size=11)

    footer(sl)


def slide_thank_you(prs):
    """Slide 17 — Thank You"""
    sl = blank(prs)
    box(sl, 0, 0, W, H, fill=NAVY)
    box(sl, 0, Inches(5.3), W, Inches(0.08), fill=WHITE)
    box(sl, 0, Inches(5.38), W, Inches(2.12), fill=RGBColor(0x0D,0x16,0x5E))

    txt(sl, "Thank You",
        Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.4),
        size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "AI-Based Intelligent Construction Planning & Management System",
        Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.6),
        size=18, color=RGBColor(0xBB,0xDE,0xFB), align=PP_ALIGN.CENTER)

    txt(sl, "R26-IT-117  |  SLIIT Faculty of Computing  |  May 2026",
        Inches(1.0), Inches(3.75), Inches(11.3), Inches(0.4),
        size=13, color=RGBColor(0x90,0xCA,0xF9), align=PP_ALIGN.CENTER)

    txt(sl, "Shazni M.I.A  |  N.H. Afaaf Hariri  |  Hanfi A.M.M  |  Moorthy T",
        Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.4),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Supervisor: Ms. Jenny Krishara  |  Co-Supervisor: Dr. Dinuka Wijendra",
        Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.4),
        size=12, color=RGBColor(0x90,0xCA,0xF9), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_title(prs)
    slide_team(prs)
    slide_agenda(prs)
    slide_system_overview(prs)
    slide_progress_summary(prs)
    slide_c01_intro(prs)
    slide_c01_pipeline(prs)
    slide_c01_stage1(prs)
    slide_c01_stage2(prs)
    slide_c01_stage3(prs)
    slide_c01_stage4(prs)
    slide_c01_metrics(prs)
    slide_c01_demo(prs)
    slide_c01_design_excellence(prs)
    slide_c01_challenges(prs)
    slide_next_steps(prs)
    slide_thank_you(prs)

    out = r"C:\Users\Asus\Desktop\R26-IT-117_PP1_Progress.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
